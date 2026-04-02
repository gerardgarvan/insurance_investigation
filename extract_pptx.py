"""Extract all text and table content from PPTX files for comparison."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import sys


def extract_pptx_content(filepath):
    """Extract all text, tables, and shape content from a PPTX file."""
    prs = Presentation(filepath)
    lines = []
    lines.append(f"{'='*80}")
    lines.append(f"FILE: {os.path.basename(filepath)}")
    lines.append(f"Total slides: {len(prs.slides)}")
    lines.append(f"{'='*80}")
    lines.append("")

    for slide_idx, slide in enumerate(prs.slides, 1):
        lines.append(f"{'─'*80}")
        lines.append(f"SLIDE {slide_idx}")
        lines.append(f"{'─'*80}")

        # Extract slide title if present
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            lines.append(f"  TITLE: {title_text}")
        else:
            lines.append(f"  TITLE: (none)")
        lines.append("")

        # Track shape index for ordering
        shape_count = 0
        for shape in slide.shapes:
            shape_count += 1

            # Handle tables
            if shape.has_table:
                table = shape.table
                lines.append(f"  [TABLE] (rows={len(table.rows)}, cols={len(table.columns)})")

                # Calculate column widths for alignment
                col_widths = []
                for col_idx in range(len(table.columns)):
                    max_width = 0
                    for row in table.rows:
                        cell_text = row.cells[col_idx].text.strip().replace('\n', ' | ')
                        max_width = max(max_width, len(cell_text))
                    col_widths.append(min(max_width + 2, 40))  # cap at 40

                for row_idx, row in enumerate(table.rows):
                    cells = []
                    for col_idx, cell in enumerate(row.cells):
                        cell_text = cell.text.strip().replace('\n', ' | ')
                        cells.append(cell_text.ljust(col_widths[col_idx]))
                    row_label = "HDR" if row_idx == 0 else f"R{row_idx:02d}"
                    lines.append(f"    {row_label}: {' │ '.join(cells)}")

                    # Add separator after header
                    if row_idx == 0:
                        sep = '─' * (sum(col_widths) + 3 * (len(col_widths) - 1))
                        lines.append(f"    {'─'*3}──{sep}")

                lines.append("")

            # Handle text frames (text boxes, titles, subtitles, etc.)
            elif shape.has_text_frame:
                # Skip the title shape since we already extracted it
                if shape == slide.shapes.title:
                    continue

                shape_type_name = "TEXT BOX"
                if hasattr(shape, 'shape_type'):
                    try:
                        shape_type_name = str(shape.shape_type).replace('MSO_SHAPE_TYPE.', '')
                    except:
                        pass

                text_content = shape.text_frame.text.strip()
                if text_content:
                    lines.append(f"  [{shape_type_name}]")
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            lines.append(f"    {para_text}")
                    lines.append("")

            # Handle grouped shapes
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                lines.append(f"  [GROUP SHAPE]")
                try:
                    for grp_shape in shape.shapes:
                        if grp_shape.has_text_frame:
                            grp_text = grp_shape.text_frame.text.strip()
                            if grp_text:
                                lines.append(f"    {grp_text}")
                        if grp_shape.has_table:
                            lines.append(f"    (nested table - {len(grp_shape.table.rows)} rows)")
                except Exception as e:
                    lines.append(f"    (could not read group: {e})")
                lines.append("")

        if shape_count == 0:
            lines.append("  (empty slide)")
            lines.append("")

    return '\n'.join(lines)


def main():
    base_dir = r"C:\Users\Owner\Documents\insurance_investigation"

    files_to_extract = [
        (os.path.join(base_dir, "insurance_tables_2026-03-24.pptx"),
         os.path.join(base_dir, "extracted_2026-03-24.txt")),
        (os.path.join(base_dir, "insurance_tables_2026-03-26.pptx"),
         os.path.join(base_dir, "extracted_2026-03-26.txt")),
    ]

    for pptx_path, output_path in files_to_extract:
        if not os.path.exists(pptx_path):
            print(f"ERROR: File not found: {pptx_path}")
            continue

        print(f"Extracting: {os.path.basename(pptx_path)} ...")
        content = extract_pptx_content(pptx_path)

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)

        print(f"  -> Saved to: {os.path.basename(output_path)}")
        print(f"  -> Lines: {content.count(chr(10)) + 1}")
        print()

    print("Done.")


if __name__ == "__main__":
    main()
